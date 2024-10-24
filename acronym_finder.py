"""
PowerPoint Acronym Finder

This script analyzes PowerPoint presentations to identify potential acronyms,
their definitions, and locations. It can optionally use predefined CSV files
of known acronyms and acronyms to exclude.

Requirements:
    python-pptx>=0.6.21
    pandas>=1.3.0
"""

import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional, Set

import pandas as pd
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.table import Table, _Cell
from pptx.slide import Slide
from pptx.util import Inches


class AcronymFinder:
    """Main class for finding and processing acronyms in PowerPoint presentations."""

    def __init__(
        self,
        pptx_path: str,
        known_acronyms_csv: Optional[str] = None,
        exclude_acronyms_csv: Optional[str] = None,
        log_level: int = logging.INFO,
    ) -> None:
        """
        Initialize the AcronymFinder.

        Args:
            pptx_path: Path to the PowerPoint file
            known_acronyms_csv: Optional path to CSV file with known acronyms
            exclude_acronyms_csv: Optional path to CSV file with acronyms to exclude
            log_level: Logging level (default: INFO)
        """
        self.pptx_path = Path(pptx_path)
        self.setup_logging(log_level)
        self.known_acronyms = self._load_known_acronyms(known_acronyms_csv) if known_acronyms_csv else {}
        self.exclusions = self._load_exclusions(exclude_acronyms_csv)
        self.found_acronyms: Dict[str, Dict] = {}

    def setup_logging(self, log_level: int) -> None:
        """Set up logging with timestamp-based filename."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        log_file = f"acronym_finder_{timestamp}.log"

        # Create logger
        self.logger = logging.getLogger(__name__)
        self.logger.setLevel(log_level)

        # Remove any existing handlers
        self.logger.handlers = []

        # Create handlers
        file_handler = logging.FileHandler(log_file)
        console_handler = logging.StreamHandler()

        # Create formatter
        formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')

        # Set formatter for handlers
        file_handler.setFormatter(formatter)
        console_handler.setFormatter(formatter)

        # Add handlers to logger
        self.logger.addHandler(file_handler)
        self.logger.addHandler(console_handler)

    def _load_known_acronyms(self, csv_path: str) -> Dict[str, str]:
        """
        Load known acronyms from CSV file.

        Args:
            csv_path: Path to CSV file containing acronyms and definitions

        Returns:
            Dictionary mapping acronyms to their definitions
        """
        try:
            df = pd.read_csv(csv_path)
            if 'Acronym' not in df.columns or 'Definition' not in df.columns:
                self.logger.error("Known acronyms CSV must contain 'Acronym' and 'Definition' columns.")
                return {}
            return dict(zip(df['Acronym'].str.upper(), df['Definition']))
        except Exception as e:
            self.logger.error(f"Error loading known acronyms: {e}")
            return {}

    def _load_exclusions(self, csv_path: Optional[str]) -> Set[str]:
        """
        Load additional acronyms to exclude from a CSV file.

        Args:
            csv_path: Optional path to CSV file containing acronyms to exclude

        Returns:
            Set of acronyms to exclude
        """
        # Default built-in exclusions
        exclusions = {'I', 'A', 'OK', 'ID', 'NO', 'AM', 'PM', 'THE'}

        if csv_path:
            try:
                df = pd.read_csv(csv_path)
                if 'Exclusion' not in df.columns:
                    self.logger.error("Exclusion CSV must contain 'Exclusion' column.")
                    return exclusions
                additional_exclusions = set(df['Exclusion'].dropna().astype(str).str.upper())
                exclusions.update(additional_exclusions)
                self.logger.info(f"Loaded {len(additional_exclusions)} exclusions from {csv_path}.")
            except Exception as e:
                self.logger.error(f"Error loading exclusions: {e}")
        else:
            self.logger.info("No exclusion CSV provided; using default exclusions.")

        return exclusions

    def _is_potential_acronym(self, word: str) -> bool:
        """
        Check if a word could be an acronym based on refined patterns.

        Args:
            word: Word to check

        Returns:
            Boolean indicating if the word might be an acronym
        """
        word_upper = word.upper()

        # Exclude if in the exclusion list
        if word_upper in self.exclusions:
            return False

        # Exclude acronyms that are entirely numeric
        if word.isdigit():
            return False

        # Exclude acronyms with hyphens followed by 2-5 digits (e.g., ABC-12345)
        if re.match(r'^[A-Z]+-\d{2,5}$', word_upper):
            return False

        # Exclude acronyms with multiple hyphen-separated numeric sections (e.g., 123-456)
        if re.match(r'^(\d+-)+\d+$', word_upper):
            return False

        # Define patterns to match valid acronyms
        acronym_patterns = [
            # Traditional uppercase acronyms (e.g., NASA)
            r'^[A-Z]{2,6}$',

            # Acronyms with a single digit followed by letters (e.g., 4CYC)
            r'^[0-9][A-Z]{1,5}$',

            # Acronyms with ampersand (e.g., I&T)
            r'^[A-Z]&[A-Z]$',

            # Acronyms with forward slash (e.g., L/TA)
            r'^[A-Z]+/[A-Z]+$',

            # Acronyms with hyphen but not followed by numbers (e.g., X-RAY)
            r'^[A-Z0-9]+-[A-Z0-9]+$',
        ]

        # Check if word matches any of the valid acronym patterns
        for pattern in acronym_patterns:
            if re.fullmatch(pattern, word_upper):
                return True

        return False

    def _find_potential_definition(self, text: str, acronym: str) -> Optional[str]:
        """
        Look for potential definition of an acronym in surrounding text.

        Args:
            text: Text to search for definition
            acronym: Acronym to find definition for

        Returns:
            Potential definition if found, None otherwise
        """
        # Escape special characters in acronym for regex
        escaped_acronym = re.escape(acronym)

        # Define patterns to capture definitions in various formats
        patterns = [
            # Standard format: acronym followed by definition
            rf'{escaped_acronym}\s*\(([\w\s,/-]+)\)',  # ABC (American Broadcasting Company)

            # Reverse format: definition followed by acronym
            rf'\(([\w\s,/-]+)\)\s*{escaped_acronym}',   # (American Broadcasting Company) ABC

            # Colon format
            rf'{escaped_acronym}:\s*([\w\s,/-]+)',      # ABC: American Broadcasting Company

            # Dash format
            rf'{escaped_acronym}\s*-\s*([\w\s,/-]+)',   # ABC - American Broadcasting Company

            # "stands for" format
            rf'{escaped_acronym}\s+stands\s+for\s+([\w\s,/-]+)',  # ABC stands for American Broadcasting Company

            # "means" format
            rf'{escaped_acronym}\s+means\s+([\w\s,/-]+)'  # ABC means American Broadcasting Company
        ]

        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1).strip()
        return None

    def _process_table_cell(self, cell: _Cell) -> str:
        """
        Extract text from a table cell.

        Args:
            cell: PowerPoint table cell object

        Returns:
            Text content of the cell
        """
        try:
            # Get text from the cell
            text = cell.text.strip()

            # Process any shapes within the cell (if supported by the PowerPoint version)
            if hasattr(cell, "shapes"):
                for shape in cell.shapes:
                    text += "\n" + self._extract_text_from_shape(shape)

            return text
        except Exception as e:
            self.logger.warning(f"Error extracting text from table cell: {e}")
            return ""

    def _process_table(self, table: Table) -> str:
        """
        Process a PowerPoint table and extract all text.

        Args:
            table: PowerPoint table object

        Returns:
            Concatenated text from all cells
        """
        table_text = []

        try:
            # Iterate through all rows and cells
            for row in table.rows:
                for cell in row.cells:
                    cell_text = self._process_table_cell(cell)
                    if cell_text:
                        table_text.append(cell_text)

            # Join all text with spaces
            return " ".join(table_text)
        except Exception as e:
            self.logger.warning(f"Error processing table: {e}")
            return ""

    def _extract_text_from_shape(self, shape: Shape) -> str:
        """
        Safely extract text from a PowerPoint shape.

        Args:
            shape: PowerPoint shape object

        Returns:
            Text content of the shape
        """
        try:
            # Handle tables
            if shape.has_table:
                self.logger.debug("Found table in slide")
                return self._process_table(shape.table)

            # Handle regular shapes with text
            elif hasattr(shape, "text"):
                return shape.text

            return ""
        except Exception as e:
            self.logger.warning(f"Error extracting text from shape: {e}")
            return ""

    def process_slide(self, slide: Slide, slide_number: int) -> None:
        """
        Process a single slide to find acronyms and their definitions.

        Args:
            slide: PowerPoint slide object
            slide_number: Number of the current slide
        """
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            # Modified to handle words with special characters
            words = re.findall(r'\b[\w/&-]+\b', text)

            for word in words:
                if self._is_potential_acronym(word):
                    word_upper = word.upper()
                    if word_upper not in self.found_acronyms:
                        self.found_acronyms[word_upper] = {
                            'definition': self.known_acronyms.get(word_upper),
                            'slides': set()
                        }

                    self.found_acronyms[word_upper]['slides'].add(slide_number)

                    # Look for definition if not already known
                    if not self.found_acronyms[word_upper]['definition']:
                        definition = self._find_potential_definition(text, word_upper)
                        if definition:
                            self.found_acronyms[word_upper]['definition'] = definition
                            self.logger.info(f"Found definition for {word_upper}: {definition}")

    def create_acronym_slide(self, prs: Presentation) -> None:
        """
        Create or update a new slide with acronym information.

        Args:
            prs: PowerPoint presentation object
        """
        # Define slide dimensions
        slide_width = prs.slide_width  # Typically 10 inches in EMUs
        slide_height = prs.slide_height

        # Convert inches to EMUs for positioning
        left_margin = Inches(0.5)
        right_margin = Inches(0.5)
        top_margin = Inches(1)
        bottom_margin = Inches(1)

        # Calculate usable width
        usable_width = slide_width - left_margin - right_margin

        # Add new slide with blank layout
        layout = prs.slide_layouts[5]  # Using blank layout
        slide = prs.slides.add_slide(layout)

        # Add title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = "Acronyms Found"
        else:
            # If the layout does not have a title placeholder, add one
            tx_box = slide.shapes.add_textbox(
                left_margin,
                top_margin - Inches(0.5),
                usable_width,
                Inches(0.5)
            )
            tf = tx_box.text_frame
            tf.text = "Acronyms Found"
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        # Create table
        rows = len(self.found_acronyms) + 1  # +1 for header
        cols = 3
        table_height = slide_height - top_margin - bottom_margin - Inches(1)  # Adjust as needed

        # Add table with 0.5” left and right margins
        table = slide.shapes.add_table(
            rows, cols, left_margin, top_margin, usable_width, table_height
        ).table

        # Set column widths proportionally
        # For example: Acronym (15%), Definition (60%), Slide Numbers (25%)
        table.columns[0].width = usable_width * 0.15
        table.columns[1].width = usable_width * 0.60
        table.columns[2].width = usable_width * 0.25

        # Set headers
        headers = ['Acronym', 'Definition', 'Slide Numbers']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            # Format header cells (e.g., bold)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

        # Populate table
        for i, (acronym, info) in enumerate(sorted(self.found_acronyms.items()), 1):
            table.cell(i, 0).text = acronym
            table.cell(i, 1).text = info['definition'] or 'Unknown'
            table.cell(i, 2).text = ', '.join(map(str, sorted(info['slides'])))

    def process_presentation(self) -> None:
        """Process the entire presentation to find and document acronyms."""
        try:
            prs = Presentation(self.pptx_path)

            # Process all slides
            for i, slide in enumerate(prs.slides, 1):
                self.logger.info(f"Processing slide {i}")
                self.process_slide(slide, i)

            # Create acronym slide
            self.create_acronym_slide(prs)

            # Save the presentation
            output_path = self.pptx_path.with_name(
                f"{self.pptx_path.stem}_with_acronyms{self.pptx_path.suffix}"
            )
            prs.save(output_path)
            self.logger.info(f"Saved presentation with acronyms to {output_path}")

        except Exception as e:
            self.logger.error(f"Error processing presentation: {e}")
            raise


def main():
    """Main entry point for the script."""
    import argparse

    parser = argparse.ArgumentParser(description="Find acronyms in PowerPoint presentations")
    parser.add_argument("pptx_path", help="Path to the PowerPoint file")
    parser.add_argument(
        "--known-acronyms",
        help="Path to CSV file with known acronyms and definitions",
    )
    parser.add_argument(
        "--exclude-acronyms",
        help="Path to CSV file with acronyms to exclude from detection",
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=["DEBUG", "INFO", "WARNING", "ERROR", "CRITICAL"],
        help="Set the logging level",
    )

    args = parser.parse_args()

    finder = AcronymFinder(
        args.pptx_path,
        known_acronyms_csv=args.known_acronyms,
        exclude_acronyms_csv=args.exclude_acronyms,
        log_level=getattr(logging, args.log_level.upper(), logging.INFO),
    )
    finder.process_presentation()


if __name__ == "__main__":
    main()
